#!/usr/bin/env python3
"""Build the six-slide Chinese advisor deck from validated RNA-seq results."""

from __future__ import annotations

import math
import os
import re
import zipfile
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/drugrep-presentation-mpl")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from matplotlib import font_manager
from matplotlib.patches import Ellipse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
PRESENTATION_DIR = Path(__file__).resolve().parent
ASSET_DIR = PRESENTATION_DIR / "assets"
OUTPUT = PRESENTATION_DIR / "HCC_RNAseq_downstream_advisor_20260821.pptx"

# Restrained academic palette.
BG, CARD, NAVY, TEAL = "F5F2EA", "FFFFFF", "17324D", "2E7D75"
GOLD, BLUE, BRICK = "C99738", "4E79A7", "A95A43"
TEXT, MUTED, GRID = "1F2933", "66737C", "D8DED9"
PALE_TEAL, PALE_GOLD, PALE_BLUE, PALE_RED = "E3EFEC", "F3E8CF", "E8EFF5", "F3E6E1"
FONT = "Microsoft YaHei"
GROUP_COLORS = {"NC": "#7D8790", "Sora": "#C99738", "Spi": "#4E79A7", "Teg": "#2E7D75"}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


def set_chart_style() -> None:
    font_path = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
    font_manager.fontManager.addfont(font_path)
    chart_font = font_manager.FontProperties(fname=font_path).get_name()
    plt.rcParams.update(
        {
            "font.family": chart_font,
            "font.size": 9,
            "axes.labelcolor": f"#{TEXT}",
            "axes.edgecolor": f"#{GRID}",
            "xtick.color": f"#{MUTED}",
            "ytick.color": f"#{MUTED}",
            "axes.unicode_minus": False,
            "figure.facecolor": "white",
            "axes.facecolor": "white",
        }
    )


def save_chart(fig: plt.Figure, name: str) -> Path:
    path = ASSET_DIR / name
    fig.savefig(path, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def make_pca() -> Path:
    scores = pd.read_csv(ROOT / "results/tables/pca_scores.tsv", sep="\t")
    fig, ax = plt.subplots(figsize=(6.4, 4.25))
    for group in ["NC", "Sora", "Spi", "Teg"]:
        subset = scores[scores["group"] == group]
        color = GROUP_COLORS[group]
        ax.scatter(subset["PC1"], subset["PC2"], s=72, color=color, edgecolor="white", linewidth=1.1, label=group, zorder=3)
        width = max(float(subset["PC1"].std()) * 4.0, 2.5)
        height = max(float(subset["PC2"].std()) * 4.0, 2.0)
        ax.add_patch(
            Ellipse(
                (subset["PC1"].mean(), subset["PC2"].mean()),
                width,
                height,
                facecolor=color,
                edgecolor=color,
                alpha=0.10,
                linewidth=1.2,
                zorder=1,
            )
        )
        for _, row in subset.iterrows():
            sample = str(row["sample"]).split(".")[-1]
            ax.annotate(sample, (row["PC1"], row["PC2"]), xytext=(4, 4), textcoords="offset points", fontsize=7, color="#58656D")
    ax.axhline(0, color="#E2E6E3", linewidth=0.8, zorder=0)
    ax.axvline(0, color="#E2E6E3", linewidth=0.8, zorder=0)
    ax.set_xlabel("PC1（70.8%）")
    ax.set_ylabel("PC2（16.3%）")
    ax.legend(frameon=False, ncol=4, loc="upper center", bbox_to_anchor=(0.5, 1.08), handletextpad=0.3, columnspacing=1.2)
    sns.despine(ax=ax)
    fig.tight_layout()
    return save_chart(fig, "pca_groups.png")


def make_volcanoes() -> Path:
    contrast_labels = [("SoravsNC", "Sora vs NC"), ("SpivsNC", "Spi vs NC"), ("TegvsNC", "Teg vs NC")]
    fig, axes = plt.subplots(1, 3, figsize=(9.1, 4.25), sharey=True)
    for ax, (contrast, title) in zip(axes, contrast_labels):
        frame = pd.read_csv(ROOT / f"data/processed/differential_expression/RSEM_DESeq2_{contrast}_mapped.tsv", sep="\t")
        frame["plot_y"] = -np.log10(frame["padj"].clip(lower=1e-35))
        sig = frame["padj"].notna() & (frame["padj"] < 0.05) & (frame["log2FoldChange"].abs() >= 1)
        up = sig & (frame["log2FoldChange"] > 0)
        down = sig & (frame["log2FoldChange"] < 0)
        ax.scatter(frame.loc[~sig, "log2FoldChange"], frame.loc[~sig, "plot_y"], s=4, color="#D3D8DA", alpha=0.42, linewidths=0, rasterized=True)
        ax.scatter(frame.loc[down, "log2FoldChange"], frame.loc[down, "plot_y"], s=6, color=f"#{BLUE}", alpha=0.64, linewidths=0, rasterized=True)
        ax.scatter(frame.loc[up, "log2FoldChange"], frame.loc[up, "plot_y"], s=6, color=f"#{GOLD}", alpha=0.68, linewidths=0, rasterized=True)
        ax.axvline(-1, color="#98A3AA", linestyle="--", linewidth=0.8)
        ax.axvline(1, color="#98A3AA", linestyle="--", linewidth=0.8)
        ax.axhline(-math.log10(0.05), color="#98A3AA", linestyle="--", linewidth=0.8)
        ax.set_title(title, fontsize=10.5, fontweight="medium", color=f"#{NAVY}", pad=7)
        ax.set_xlabel("log2FC")
        ax.set_xlim(-8, 8)
        sns.despine(ax=ax)
    axes[0].set_ylabel("−log10(FDR)")
    fig.text(0.50, 0.01, "蓝：下调 DEG    金：上调 DEG    灰：未达阈值", ha="center", fontsize=8, color=f"#{MUTED}")
    fig.tight_layout(rect=(0, 0.05, 1, 1), w_pad=1.4)
    return save_chart(fig, "volcano_triptych.png")


def make_robustness() -> Path:
    data = pd.read_csv(ROOT / "results/tables/robustness_summary.tsv", sep="\t")
    names = {
        "RSEM DESeq2 vs Salmon DESeq2": "RSEM–Salmon",
        "RSEM DESeq2 vs STAR DESeq2": "RSEM–STAR",
        "RSEM DESeq2 vs RSEM edgeR": "DESeq2–edgeR",
        "RSEM DESeq2 vs RSEM Limma": "DESeq2–limma",
    }
    contrasts = {"SoravsNC": "Sora", "SpivsNC": "Spi", "TegvsNC": "Teg"}
    data["method"] = data["comparison"].map(names)
    data["treatment"] = data["contrast"].map(contrasts)
    order = ["RSEM–Salmon", "RSEM–STAR", "DESeq2–edgeR", "DESeq2–limma"]
    ybase = {name: len(order) - 1 - index for index, name in enumerate(order)}
    offsets = {"Sora": 0.18, "Spi": 0.0, "Teg": -0.18}
    fig, ax = plt.subplots(figsize=(6.0, 3.75))
    ax.axvspan(0.90, 1.005, color="#E7F0ED", zorder=0)
    for method in order:
        subset = data[data["method"] == method]
        ax.hlines(ybase[method], subset["log2fc_spearman"].min(), subset["log2fc_spearman"].max(), color="#BCC5C7", linewidth=1.2, zorder=1)
    for treatment in ["Sora", "Spi", "Teg"]:
        subset = data[data["treatment"] == treatment]
        ax.scatter(
            subset["log2fc_spearman"],
            [ybase[row] + offsets[treatment] for row in subset["method"]],
            s=48,
            color=GROUP_COLORS[treatment],
            edgecolor="white",
            linewidth=0.8,
            label=treatment,
            zorder=3,
        )
    ax.axvline(0.90, color=f"#{TEAL}", linewidth=0.9, linestyle="--")
    ax.set_xlim(0.75, 1.01)
    ax.set_xticks([0.75, 0.80, 0.85, 0.90, 0.95, 1.00])
    ax.set_xlabel("log2FC Spearman ρ")
    ax.set_yticks(list(ybase.values()), list(ybase.keys()))
    ax.legend(frameon=False, ncol=3, loc="lower center", bbox_to_anchor=(0.5, 1.01), columnspacing=1.0, handletextpad=0.3)
    ax.grid(axis="x", color="#E5E8E5", linewidth=0.7)
    sns.despine(ax=ax, left=True)
    fig.tight_layout()
    return save_chart(fig, "method_robustness.png")


def make_similarity() -> Path:
    treatment = pd.read_csv(ROOT / "results/tables/treatment_similarity.tsv", sep="\t")
    pathway = pd.read_csv(ROOT / "results/tables/pathway_similarity.tsv", sep="\t")

    def pick(frame: pd.DataFrame, a: str, b: str, column: str) -> float:
        row = frame[((frame["contrast_a"] == a) & (frame["contrast_b"] == b)) | ((frame["contrast_a"] == b) & (frame["contrast_b"] == a))]
        return float(row.iloc[0][column])

    metrics = ["全基因统计量 ρ", "通路 NES ρ", "DEG Jaccard"]
    spi = [pick(treatment, "SoravsNC", "SpivsNC", "stat_spearman"), pick(pathway, "SoravsNC", "SpivsNC", "nes_spearman"), pick(treatment, "SoravsNC", "SpivsNC", "deg_jaccard")]
    teg = [pick(treatment, "SoravsNC", "TegvsNC", "stat_spearman"), pick(pathway, "SoravsNC", "TegvsNC", "nes_spearman"), pick(treatment, "SoravsNC", "TegvsNC", "deg_jaccard")]
    y = np.arange(len(metrics))
    fig, ax = plt.subplots(figsize=(5.6, 3.75))
    h = 0.28
    ax.barh(y + h / 2, spi, height=h, color="#AAB2B7", label="Sora–Spi")
    ax.barh(y - h / 2, teg, height=h, color=f"#{TEAL}", label="Sora–Teg")
    for values, offset in [(spi, h / 2), (teg, -h / 2)]:
        for index, value in enumerate(values):
            ax.text(value + 0.018, index + offset, f"{value:.3f}", va="center", fontsize=8, color=f"#{TEXT}")
    ax.set_yticks(y, metrics)
    ax.invert_yaxis()
    ax.set_xlim(0, 0.82)
    ax.set_xlabel("相似度（描述性）")
    ax.legend(frameon=False, ncol=2, loc="lower center", bbox_to_anchor=(0.5, 1.02), columnspacing=1.2)
    ax.grid(axis="x", color="#E5E8E5", linewidth=0.7)
    sns.despine(ax=ax, left=True)
    fig.tight_layout()
    return save_chart(fig, "sorafenib_similarity.png")


def make_pathway_heatmap() -> Path:
    gsea = pd.read_csv(ROOT / "results/tables/gsea_all.tsv", sep="\t")
    gsea = gsea[(gsea["quantifier"] == "RSEM") & (gsea["library"] == "MSigDB_Hallmark_2020")].copy()
    selected = [
        ("Myc Targets V1", "MYC 靶基因 V1"),
        ("E2F Targets", "E2F 靶基因"),
        ("G2-M Checkpoint", "G2/M 检查点"),
        ("Unfolded Protein Response", "未折叠蛋白反应"),
        ("TNF-alpha Signaling via NF-kB", "TNFα / NF-κB"),
        ("IL-6/JAK/STAT3 Signaling", "IL-6 / JAK / STAT3"),
        ("mTORC1 Signaling", "mTORC1"),
        ("Hypoxia", "缺氧"),
        ("Glycolysis", "糖酵解"),
        ("Coagulation", "凝血"),
    ]
    contrasts = ["SoravsNC", "SpivsNC", "TegvsNC"]
    matrix = np.zeros((len(selected), len(contrasts)))
    fdr = np.ones_like(matrix)
    for row_index, (term, _) in enumerate(selected):
        for column_index, contrast in enumerate(contrasts):
            record = gsea[(gsea["term"] == term) & (gsea["contrast"] == contrast)].iloc[0]
            matrix[row_index, column_index] = record["nes"]
            fdr[row_index, column_index] = record["fdr"]
    labels = [[f"{matrix[i, j]:.1f}{'*' if fdr[i, j] < 0.05 else ''}" for j in range(3)] for i in range(len(selected))]
    fig, ax = plt.subplots(figsize=(7.0, 5.15))
    sns.heatmap(
        matrix,
        ax=ax,
        cmap=sns.diverging_palette(240, 20, as_cmap=True),
        center=0,
        vmin=-2.8,
        vmax=2.8,
        annot=np.array(labels),
        fmt="",
        annot_kws={"fontsize": 8.4},
        linewidths=1.0,
        linecolor="white",
        cbar_kws={"label": "NES", "shrink": 0.72},
    )
    ax.set_xticklabels(["Sora", "Spi", "Teg"], rotation=0, fontsize=9)
    ax.set_yticklabels([label for _, label in selected], rotation=0, fontsize=8.5)
    ax.set_xlabel("处理组 vs NC")
    ax.set_ylabel("")
    ax.tick_params(length=0)
    fig.text(0.01, 0.01, "* FDR < 0.05；正/负 NES 表示在排序两端富集，不直接等同于通路激活/抑制", fontsize=7.5, color=f"#{MUTED}")
    fig.tight_layout(rect=(0, 0.04, 1, 1))
    return save_chart(fig, "hallmark_heatmap.png")


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: float = 14, color: str = TEXT, bold: bool = False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin: float = 0.02, font: str = FONT, line_spacing: float = 1.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_paragraphs(slide, items: list[dict], x: float, y: float, w: float, h: float, margin: float = 0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_after = Pt(item.get("space_after", 5))
        paragraph.line_spacing = item.get("line_spacing", 1.04)
        run = paragraph.add_run()
        run.text = item["text"]
        run.font.name = item.get("font", FONT)
        run.font.size = Pt(item.get("size", 13))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = rgb(item.get("color", TEXT))
    return box


def add_card(slide, x: float, y: float, w: float, h: float, fill: str = CARD, line: str = "E2E5E2", radius: bool = True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.7)
    return shape


def add_rule(slide, x: float, y: float, w: float, color: str = GOLD, height: float = 0.035):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_slide_header(slide, number: int, title: str, kicker: str | None = None):
    add_text(slide, f"0{number}", 0.48, 0.22, 0.52, 0.30, size=9, color=GOLD, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "DRUGREP · HCC RNA-SEQ", 1.00, 0.22, 2.6, 0.30, size=8.5, color=MUTED, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, 0.48, 0.65, 10.9, 0.52, size=24, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if kicker:
        add_text(slide, kicker, 10.9, 0.75, 1.95, 0.30, size=9, color=TEAL, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
    add_rule(slide, 0.48, 1.26, 12.35, color=GRID, height=0.012)


def add_footer(slide, number: int, note: str = "内部讨论 · 2026-08-21"):
    add_text(slide, note, 0.50, 7.18, 4.5, 0.18, size=7.5, color=MUTED)
    add_text(slide, str(number), 12.48, 7.16, 0.35, 0.20, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    from PIL import Image

    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))


def make_deck(assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(BG)
        return slide

    # 1 — title
    slide = new_slide()
    add_rule(slide, 0.55, 0.55, 1.05, color=GOLD, height=0.07)
    add_text(slide, "DRUGREP · TRANSCRIPTOMICS", 0.55, 0.77, 3.6, 0.32, size=9, color=MUTED, bold=True)
    add_text(slide, "HCC 候选药物", 0.55, 1.48, 6.5, 0.64, size=31, color=NAVY, bold=True)
    add_text(slide, "RNA-seq 下游分析", 0.55, 2.18, 6.5, 0.66, size=31, color=TEAL, bold=True)
    add_text(slide, "对 ZZH 前期核心表格的复核、稳健性分析与通路线索", 0.58, 3.02, 6.9, 0.55, size=15, color=TEXT)
    add_card(slide, 8.10, 1.15, 4.50, 4.45, fill=CARD, line="E0E4E0")
    add_text(slide, "从样本到可讨论的机制线索", 8.45, 1.48, 3.85, 0.40, size=13.5, color=NAVY, bold=True)
    y_positions = [2.10, 2.82, 3.54, 4.26]
    for (group, fill), y in zip([("NC", "ECEFED"), ("Sora", PALE_GOLD), ("Spi", PALE_BLUE), ("Teg", PALE_TEAL)], y_positions):
        add_card(slide, 8.45, y, 1.16, 0.48, fill=fill, line=fill)
        add_text(slide, f"{group}  ×3", 8.45, y, 1.16, 0.48, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.68), Inches(y + 0.24), Inches(10.35), Inches(3.42))
        connector.line.color.rgb = rgb("B7C1BE")
        connector.line.width = Pt(1.2)
    add_card(slide, 10.35, 2.72, 1.82, 1.40, fill=PALE_TEAL, line="B9D1CB")
    add_text(slide, "QC → DEG\n→ GSEA / ORA", 10.48, 2.84, 1.56, 1.12, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    for x, label, fill in [(0.58, "12 个样本", PALE_BLUE), (2.16, "3 个处理", PALE_GOLD), (3.74, "3 套定量", PALE_TEAL)]:
        add_card(slide, x, 4.37, 1.38, 0.50, fill=fill, line=fill)
        add_text(slide, label, x, 4.37, 1.38, 0.50, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "结论边界：本次仅比较药物扰动，不能据此宣称 HCC 逆转或疗效。", 0.58, 5.43, 6.9, 0.48, size=11, color=BRICK)
    add_text(slide, "2026-08-21", 0.58, 6.60, 2.0, 0.28, size=9, color=MUTED)
    add_text(slide, "RNA-seq downstream review", 9.2, 6.60, 3.4, 0.28, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    # 2 — outline and workflow
    slide = new_slide()
    add_slide_header(slide, 2, "汇报结构与下游分析路径", "先验证，再解释")
    add_text(slide, "汇报结构", 0.58, 1.55, 2.7, 0.35, size=13, color=TEAL, bold=True)
    outline = [
        ("01", "数据与质量", "输入是否完整、样本是否可信"),
        ("02", "差异表达", "三种处理的响应幅度与方向"),
        ("03", "稳健性", "定量方法与统计方法是否一致"),
        ("04", "通路与判断", "候选药物差异、边界、下一步"),
    ]
    for index, (num, label, desc) in enumerate(outline):
        y = 2.00 + index * 1.03
        add_text(slide, num, 0.60, y, 0.48, 0.40, size=10, color=GOLD, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, 1.13, y, 1.48, 0.40, size=14, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, 1.13, y + 0.39, 2.45, 0.38, size=9.5, color=MUTED)
        if index < 3:
            add_rule(slide, 1.13, y + 0.87, 2.45, color="DDE2DF", height=0.012)
    add_card(slide, 4.12, 1.55, 8.58, 4.98, fill=CARD, line="E1E5E2")
    add_text(slide, "已完成的下游分析步骤", 4.48, 1.88, 4.2, 0.35, size=13, color=NAVY, bold=True)
    steps = [
        ("1", "输入核对", "校验和 / 表结构 / 样本对应"),
        ("2", "表达 QC", "库大小 / 相关性 / PCA / 热图"),
        ("3", "DEG", "RSEM–DESeq2 为主；复现阈值"),
        ("4", "稳健性", "Salmon / STAR；edgeR / limma"),
        ("5", "通路", "全基因 Wald 排序 GSEA；上下调 ORA"),
        ("6", "横向比较", "基因、DEG、NES、leading edge"),
    ]
    for index, (num, label, desc) in enumerate(steps):
        row, col = divmod(index, 3)
        x = 4.48 + col * 2.58
        y = 2.52 + row * 1.72
        fill = [PALE_BLUE, PALE_GOLD, PALE_TEAL, "EDF0F1", PALE_GOLD, PALE_TEAL][index]
        add_card(slide, x, y, 2.22, 1.28, fill=fill, line=fill)
        add_text(slide, num, x + 0.16, y + 0.13, 0.30, 0.28, size=9, color=GOLD, bold=True)
        add_text(slide, label, x + 0.48, y + 0.12, 1.52, 0.30, size=12, color=NAVY, bold=True)
        add_text(slide, desc, x + 0.16, y + 0.53, 1.88, 0.54, size=9.2, color=TEXT)
    add_footer(slide, 2)

    # 3 — raw data and QC
    slide = new_slide()
    add_slide_header(slide, 3, "原始数据与质量控制", "12 样本 · 4 组 · n=3")
    add_text(slide, "输入组成", 0.58, 1.52, 2.1, 0.32, size=13, color=TEAL, bold=True)
    for index, group in enumerate(["NC", "Sora", "Spi", "Teg"]):
        x = 0.58 + index * 1.18
        fill = {"NC": "ECEFED", "Sora": PALE_GOLD, "Spi": PALE_BLUE, "Teg": PALE_TEAL}[group]
        add_card(slide, x, 1.93, 0.98, 0.67, fill=fill, line=fill)
        add_text(slide, group, x, 1.99, 0.98, 0.25, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "3 replicates", x, 2.28, 0.98, 0.20, size=7.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 0.58, 2.86, 4.55, 2.75, fill=CARD, line="E1E5E2")
    add_paragraphs(
        slide,
        [
            {"text": "ZZH 已完成", "size": 11, "bold": True, "color": NAVY, "space_after": 3},
            {"text": "QC / 比对 / 表达定量 / DESeq2 / edgeR / limma", "size": 10.2, "color": TEXT, "space_after": 10},
            {"text": "本次复核", "size": 11, "bold": True, "color": NAVY, "space_after": 3},
            {"text": "52 个校验和；样本名一致；计数非负；DEG 阈值完全复现", "size": 10.2, "color": TEXT, "space_after": 10},
            {"text": "仍缺失", "size": 11, "bold": True, "color": BRICK, "space_after": 3},
            {"text": "FASTQ 级 QC、比对汇总、完整实验与批次信息", "size": 10.2, "color": TEXT, "space_after": 0},
        ],
        0.82,
        3.10,
        4.06,
        2.30,
    )
    add_card(slide, 0.58, 5.82, 4.55, 0.68, fill=PALE_RED, line="E4C7BD")
    add_text(slide, "注释合并最后一步报错，但核心计数与 DE 表均可用。", 0.80, 5.95, 4.10, 0.34, size=10, color=BRICK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_card(slide, 5.45, 1.53, 7.27, 4.98, fill=CARD, line="E1E5E2")
    add_text(slide, "PCA：各组分离，组内重复紧密", 5.80, 1.82, 3.8, 0.34, size=13, color=NAVY, bold=True)
    add_picture_contain(slide, assets["pca"], 5.75, 2.20, 6.60, 3.70)
    add_text(slide, "PC1 70.8% · PC2 16.3%", 5.84, 6.03, 2.7, 0.24, size=9, color=MUTED)
    add_text(slide, "组内平均相关性 0.9967–0.9981；未标记离群样本", 8.33, 6.03, 3.95, 0.24, size=9, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 3)

    # 4 — differential expression
    slide = new_slide()
    add_slide_header(slide, 4, "差异表达：Sora 响应最广，Teg 次之，Spi 较窄", "RSEM–DESeq2")
    metrics = [
        ("Sora vs NC", "2,246", "1,292 ↑  /  954 ↓", PALE_GOLD),
        ("Spi vs NC", "421", "192 ↑  /  229 ↓", PALE_BLUE),
        ("Teg vs NC", "657", "442 ↑  /  215 ↓", PALE_TEAL),
    ]
    for index, (label, total, split, fill) in enumerate(metrics):
        x = 0.58 + index * 2.06
        add_card(slide, x, 1.53, 1.80, 1.23, fill=fill, line=fill)
        add_text(slide, label, x + 0.14, 1.68, 1.52, 0.26, size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, total, x + 0.14, 1.97, 1.52, 0.42, size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, split, x + 0.14, 2.43, 1.52, 0.20, size=8, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 6.90, 1.53, 5.82, 1.23, fill=CARD, line="E1E5E2")
    add_text(slide, "判定阈值", 7.20, 1.75, 1.10, 0.25, size=10, color=MUTED, bold=True)
    add_text(slide, "FDR < 0.05  且  |log2FC| ≥ 1", 7.20, 2.04, 3.25, 0.38, size=16, color=NAVY, bold=True)
    add_text(slide, "三套定量的存档 DEG 均被精确复现", 9.45, 2.04, 2.80, 0.38, size=10, color=TEAL, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
    add_card(slide, 0.58, 3.03, 12.14, 3.55, fill=CARD, line="E1E5E2")
    add_picture_contain(slide, assets["volcano"], 0.87, 3.24, 11.56, 2.92)
    add_text(slide, "DEG 数量反映扰动幅度，不直接代表药物优劣。", 0.92, 6.22, 5.4, 0.24, size=9, color=BRICK)
    add_text(slide, "Sora 的 DEG 数约为 Teg 的 3.4×、Spi 的 5.3×", 7.05, 6.22, 5.1, 0.24, size=9, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 4)

    # 5 — robustness and similarity
    slide = new_slide()
    add_slide_header(slide, 5, "稳健性与候选药物相似性", "方向一致性优先")
    add_card(slide, 0.58, 1.52, 6.15, 4.82, fill=CARD, line="E1E5E2")
    add_text(slide, "跨定量 / 统计方法的 log2FC 一致性", 0.90, 1.82, 4.5, 0.34, size=12.5, color=NAVY, bold=True)
    add_picture_contain(slide, assets["robustness"], 0.90, 2.18, 5.55, 3.55)
    add_text(slide, "共享 DEG 的方向一致率：100%", 0.92, 5.84, 2.9, 0.25, size=9.5, color=TEAL, bold=True)
    add_text(slide, "Hallmark NES：RSEM vs Salmon/STAR ρ=0.989–0.999", 2.95, 5.84, 3.42, 0.25, size=8.8, color=MUTED, align=PP_ALIGN.RIGHT)
    add_card(slide, 6.98, 1.52, 5.74, 4.82, fill=CARD, line="E1E5E2")
    add_text(slide, "相对 Sora 的扰动相似性", 7.30, 1.82, 3.6, 0.34, size=12.5, color=NAVY, bold=True)
    add_picture_contain(slide, assets["similarity"], 7.28, 2.20, 5.05, 3.25)
    add_card(slide, 7.34, 5.51, 4.98, 0.56, fill=PALE_TEAL, line=PALE_TEAL)
    add_text(slide, "Teg 与 Sora 更接近：ρgene=0.568；ρpathway=0.722", 7.52, 5.63, 4.60, 0.29, size=9.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_card(slide, 0.58, 6.55, 12.14, 0.40, fill=PALE_RED, line=PALE_RED)
    add_text(slide, "相似性仅描述转录组扰动，不是疗效分数；DEG 重叠也不能替代全基因与通路比较。", 0.80, 6.61, 11.70, 0.25, size=9.5, color=BRICK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 5)

    # 6 — pathways and next steps
    slide = new_slide()
    add_slide_header(slide, 6, "通路线索、判断边界与下一步", "与 ZZH 对齐后推进")
    add_card(slide, 0.58, 1.52, 7.22, 4.92, fill=CARD, line="E1E5E2")
    add_text(slide, "Hallmark GSEA：候选药物呈现不同程序", 0.90, 1.80, 4.3, 0.34, size=12.5, color=NAVY, bold=True)
    add_picture_contain(slide, assets["pathway"], 0.86, 2.18, 6.62, 3.93)
    add_text(slide, "ORA 印证：三组 TNF 相关下调；Sora 凝血/补体上调；Teg 缺氧相关上调", 0.92, 6.09, 6.48, 0.24, size=8.7, color=MUTED)
    add_card(slide, 8.05, 1.52, 4.67, 2.37, fill=CARD, line="E1E5E2")
    add_paragraphs(
        slide,
        [
            {"text": "当前判断", "size": 12.5, "bold": True, "color": NAVY, "space_after": 7},
            {"text": "Teg：MYC / E2F / G2M 负向富集；伴随缺氧、糖酵解正向富集", "size": 9.6, "color": TEXT, "space_after": 7},
            {"text": "Spi：更偏 TNF/NF-κB、IL/JAK/STAT 与应激程序", "size": 9.6, "color": TEXT, "space_after": 7},
            {"text": "优先区分：Teg 的 Sora-like 收敛 vs Spi 的差异机制", "size": 9.6, "bold": True, "color": TEAL, "space_after": 0},
        ],
        8.34,
        1.80,
        4.08,
        1.84,
    )
    add_card(slide, 8.05, 4.10, 4.67, 1.52, fill=PALE_GOLD, line="E4D2AA")
    add_paragraphs(
        slide,
        [
            {"text": "需要 ZZH 确认", "size": 11.5, "bold": True, "color": NAVY, "space_after": 5},
            {"text": "• Spi / Teg 全名；细胞系、剂量、处理时长", "size": 9.3, "color": TEXT, "space_after": 3},
            {"text": "• 基因组/注释版本、建库与链特异性", "size": 9.3, "color": TEXT, "space_after": 3},
            {"text": "• 批次设计；FastQC/MultiQC 与比对汇总", "size": 9.3, "color": TEXT, "space_after": 0},
        ],
        8.32,
        4.31,
        4.10,
        1.10,
    )
    add_card(slide, 8.05, 5.84, 4.67, 0.60, fill=PALE_RED, line="E4C7BD")
    add_text(slide, "无肿瘤–正常签名：不能建立 HCC 逆转或疗效结论", 8.26, 5.96, 4.25, 0.32, size=9.5, color=BRICK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "建议下一步：用 qPCR / 蛋白与增殖、凋亡、细胞周期表型验证 leading-edge 基因。", 0.80, 6.68, 11.75, 0.31, size=10.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6, note="内部讨论 · 结论需结合实验元数据")

    prs.save(OUTPUT)


def validate_deck() -> pd.DataFrame:
    with zipfile.ZipFile(OUTPUT) as archive:
        bad = archive.testzip()
        if bad:
            raise RuntimeError(f"Corrupt PPTX member: {bad}")
        media = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        if len(media) < 5:
            raise RuntimeError(f"Expected at least five embedded chart images, found {len(media)}")

    deck = Presentation(OUTPUT)
    if len(deck.slides) != 6:
        raise RuntimeError(f"Expected 6 slides, found {len(deck.slides)}")
    audit = []
    for index, slide in enumerate(deck.slides, start=1):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.has_text_frame)
        if re.search(r"\bmentor\b|导师", text, re.IGNORECASE):
            raise RuntimeError(f"Slide {index} refers to the mentor without using ZZH")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > deck.slide_width or shape.top + shape.height > deck.slide_height:
                raise RuntimeError(f"Slide {index} contains an out-of-bounds shape: {shape.name}")
        chinese_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
        audit.append(
            {
                "slide": index,
                "characters": len(text),
                "chinese_characters": chinese_chars,
                "text_boxes": sum(1 for shape in slide.shapes if hasattr(shape, "text_frame") and shape.has_text_frame),
                "pictures": sum(1 for shape in slide.shapes if shape.shape_type == 13),
            }
        )
    frame = pd.DataFrame(audit)
    if frame["characters"].max() > 430:
        raise RuntimeError("A slide exceeds the text-density guardrail")
    frame.to_csv(PRESENTATION_DIR / "slide_text_audit.tsv", sep="\t", index=False)
    return frame


def main() -> int:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    set_chart_style()
    assets = {
        "pca": make_pca(),
        "volcano": make_volcanoes(),
        "robustness": make_robustness(),
        "similarity": make_similarity(),
        "pathway": make_pathway_heatmap(),
    }
    make_deck(assets)
    audit = validate_deck()
    print(f"Created {OUTPUT}")
    print(audit.to_string(index=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
